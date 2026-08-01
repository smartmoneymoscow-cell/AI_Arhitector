"""
shared/pptx_generator.py — Генератор инвестиционных презентаций (PPTX/PDF).

Генерирует:
    - Презентации для инвесторов
    - Концептуальные презентации
    - Отчёты по проекту

Использование:
    from shared.pptx_generator import PresentationGenerator

    gen = PresentationGenerator()
    path = gen.generate(data, output_path="presentation.pptx")
"""

import logging
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class Slide:
    """Один слайд презентации."""

    title: str
    content: list[str] = field(default_factory=list)
    layout: str = "title_and_content"  # title, title_and_content, two_column, image, chart
    notes: str = ""


@dataclass
class Presentation:
    """Презентация."""

    title: str
    subtitle: str = ""
    author: str = ""
    slides: list[Slide] = field(default_factory=list)

    def add_slide(self, slide: Slide):
        self.slides.append(slide)


class PresentationGenerator:
    """Генератор презентаций."""

    def generate(self, data: dict, output_path: str = "presentation.pptx") -> str:
        """
        Генерация PPTX-презентации.

        Args:
            data: данные проекта
            output_path: путь для сохранения

        Returns:
            путь к файлу
        """
        try:
            from pptx import Presentation as PptxPresentation

            pass  # PP_ALIGN unused
            pass  # Inches, Pt unused
        except ImportError:
            logger.warning("python-pptx not installed — generating HTML fallback")
            return self._generate_html(data, output_path.replace(".pptx", ".html"))

        prs = PptxPresentation()

        # Титульный слайд
        self._add_title_slide(prs, data)

        # Концепция
        if data.get("concept"):
            self._add_concept_slide(prs, data)

        # Архитектура
        if data.get("architecture"):
            self._add_architecture_slide(prs, data)

        # Мастер-план
        if data.get("masterplan"):
            self._add_masterplan_slide(prs, data)

        # Этапы реализации
        if data.get("phasing"):
            self._add_phasing_slide(prs, data)

        # Финансовая модель
        if data.get("financial"):
            self._add_financial_slide(prs, data)

        # Окупаемость
        if data.get("payback"):
            self._add_payback_slide(prs, data)

        # Капитализация
        if data.get("capitalization"):
            self._add_capitalization_slide(prs, data)

        # Инвестиционная стратегия
        if data.get("investment_strategy"):
            self._add_strategy_slide(prs, data)

        # Дорожная карта
        if data.get("roadmap"):
            self._add_roadmap_slide(prs, data)

        # Контакты
        self._add_contacts_slide(prs, data)

        prs.save(output_path)
        logger.info(f"Presentation saved: {output_path}")
        return output_path

    def _add_title_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = data.get("title", "Архитектурный проект")
            if slide.placeholders[1]:
                slide.placeholders[1].text = data.get("subtitle", "")
        except Exception:
            pass

    def _add_concept_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Концепция проекта"
            body = slide.placeholders[1]
            concept = data.get("concept", {})
            body.text = concept.get("description", "")
        except Exception:
            pass

    def _add_architecture_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Архитектурная концепция"
            body = slide.placeholders[1]
            arch = data.get("architecture", {})
            points = arch.get("key_points", [])
            body.text = "\n".join(f"• {p}" for p in points) if points else "Архитектурная концепция проекта"
        except Exception:
            pass

    def _add_masterplan_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Мастер-план"
            body = slide.placeholders[1]
            mp = data.get("masterplan", {})
            body.text = mp.get("description", "Генеральный план развития территории")
        except Exception:
            pass

    def _add_phasing_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Этапы реализации"
            body = slide.placeholders[1]
            phases = data.get("phasing", [])
            lines = []
            for i, phase in enumerate(phases, 1):
                if isinstance(phase, dict):
                    lines.append(f"Фаза {i}: {phase.get('name', '')} — {phase.get('duration', '')}")
                else:
                    lines.append(f"Фаза {i}: {phase}")
            body.text = "\n".join(lines) if lines else "Этапы реализации проекта"
        except Exception:
            pass

    def _add_financial_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Финансовая модель"
            body = slide.placeholders[1]
            fin = data.get("financial", {})
            lines = []
            if "capex" in fin:
                lines.append(f"Инвестиции: {fin['capex']:,.0f} ₽")
            if "revenue" in fin:
                lines.append(f"Выручка (год): {fin['revenue']:,.0f} ₽")
            if "ebitda" in fin:
                lines.append(f"EBITDA: {fin['ebitda']:,.0f} ₽")
            if "npv" in fin:
                lines.append(f"NPV: {fin['npv']:,.0f} ₽")
            if "irr" in fin:
                lines.append(f"IRR: {fin['irr']}%")
            body.text = "\n".join(lines) if lines else "Финансовые показатели проекта"
        except Exception:
            pass

    def _add_payback_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Окупаемость"
            body = slide.placeholders[1]
            pb = data.get("payback", {})
            lines = []
            if "years" in pb:
                lines.append(f"Срок окупаемости: {pb['years']} лет")
            if "scenarios" in pb:
                for name, years in pb["scenarios"].items():
                    lines.append(f"  {name}: {years} лет")
            body.text = "\n".join(lines) if lines else "Анализ окупаемости"
        except Exception:
            pass

    def _add_capitalization_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Капитализация"
            body = slide.placeholders[1]
            cap = data.get("capitalization", {})
            lines = []
            if "current_value" in cap:
                lines.append(f"Текущая стоимость: {cap['current_value']:,.0f} ₽")
            if "projected_value" in cap:
                lines.append(f"Прогнозная стоимость: {cap['projected_value']:,.0f} ₽")
            if "growth_pct" in cap:
                lines.append(f"Рост стоимости: {cap['growth_pct']}%")
            body.text = "\n".join(lines) if lines else "Оценка капитализации"
        except Exception:
            pass

    def _add_strategy_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Инвестиционная стратегия"
            body = slide.placeholders[1]
            strategy = data.get("investment_strategy", {})
            body.text = strategy.get("description", "Стратегия привлечения инвестиций")
        except Exception:
            pass

    def _add_roadmap_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Дорожная карта"
            body = slide.placeholders[1]
            roadmap = data.get("roadmap", [])
            lines = []
            for item in roadmap:
                if isinstance(item, dict):
                    lines.append(
                        f"• {item.get('phase', '')}: {item.get('description', '')} ({item.get('duration', '')})"
                    )
                else:
                    lines.append(f"• {item}")
            body.text = "\n".join(lines) if lines else "Дорожная карта реализации"
        except Exception:
            pass

    def _add_contacts_slide(self, prs, data: dict):
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Контакты"
            if slide.placeholders[1]:
                slide.placeholders[1].text = data.get("contacts", "Спасибо за внимание!")
        except Exception:
            pass

    def _generate_html(self, data: dict, output_path: str) -> str:
        """HTML-фолбэк если python-pptx не установлен."""
        slides_html = []

        # Title
        slides_html.append(f"""
        <div class="slide title-slide">
            <h1>{data.get("title", "Архитектурный проект")}</h1>
            <p>{data.get("subtitle", "")}</p>
        </div>""")

        # Financial
        fin = data.get("financial", {})
        if fin:
            items = "".join(
                f"<li><strong>{k}:</strong> {v:,.0f} ₽</li>" for k, v in fin.items() if isinstance(v, int | float)
            )
            slides_html.append(f"""
        <div class="slide">
            <h2>Финансовая модель</h2>
            <ul>{items}</ul>
        </div>""")

        # Contacts
        slides_html.append(f"""
        <div class="slide">
            <h2>Контакты</h2>
            <p>{data.get("contacts", "Спасибо за внимание!")}</p>
        </div>""")

        html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{data.get("title", "Презентация")}</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 0; }}
.slide {{ width: 960px; height: 540px; margin: 20px auto; padding: 40px; box-sizing: border-box;
          border: 1px solid #ddd; page-break-after: always; display: flex; flex-direction: column; justify-content: center; }}
.title-slide {{ background: #1a1a2e; color: white; text-align: center; }}
h1 {{ font-size: 36px; }} h2 {{ font-size: 28px; color: #1a1a2e; }}
ul {{ font-size: 18px; }} li {{ margin: 10px 0; }}
</style></head><body>{"".join(slides_html)}</body></html>"""

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        return output_path


# Глобальный экземпляр
_gen: PresentationGenerator | None = None


def get_presentation_generator() -> PresentationGenerator:
    global _gen
    if _gen is None:
        _gen = PresentationGenerator()
    return _gen
